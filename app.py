import streamlit as st
import arxiv
import pandas as pd
from wordcloud import WordCloud
import matplotlib.pyplot as plt
from textblob import TextBlob
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import nltk
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize, sent_tokenize
import heapq
import io

# NLTK Data for AI Summary
try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt')
    nltk.download('stopwords')

# Graphviz safe import
try:
    import graphviz
    GRAPHVIZ_AVAILABLE = True
except ImportError:
    GRAPHVIZ_AVAILABLE = False

# --- PAGE CONFIG ---
st.set_page_config(page_title="Suga's Research AI", layout="wide")

st.title("🔬 Suga’s AI: The Research Bridge")
st.markdown("### Connecting PhD Research to Young Minds 🎓")

# --- UTILITY FUNCTIONS ---

def extract_text_from_file(uploaded_file):
    file_extension = uploaded_file.name.split('.')[-1].lower()
    text = ""
    if file_extension == 'docx':
        doc = Document(uploaded_file)
        text = "\n".join([p.text for p in doc.paragraphs])
    elif file_extension == 'pptx':
        prs = Presentation(uploaded_file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"): text += shape.text + " "
    elif file_extension == 'pdf':
        reader = PdfReader(uploaded_file)
        for page in reader.pages: text += page.extract_text()
    return text

def generate_ai_summary(text, num_sentences=2):
    if len(text) < 150: return text
    stop_words = set(stopwords.words('english'))
    words = word_tokenize(text.lower())
    freq = {}
    for word in words:
        if word not in stop_words and word.isalnum():
            freq[word] = freq.get(word, 0) + 1
    
    sentences = sent_tokenize(text)
    scores = {}
    for sent in sentences:
        for word in word_tokenize(sent.lower()):
            if word in freq:
                if len(sent.split(' ')) < 35:
                    scores[sent] = scores.get(sent, 0) + freq[word]
    
    summary = heapq.nlargest(num_sentences, scores, key=scores.get)
    return " ".join(summary)

def get_career_advice(q):
    q = q.lower()
    if 'ai' in q or 'learning' in q: return ["Python", "Linear Algebra", "Calculus", "Neural Networks"]
    if 'space' in q or 'physic' in q: return ["Advanced Physics", "Calculus", "Astrophysics", "Data Science"]
    if 'bio' in q or 'med' in q: return ["Organic Chemistry", "Biology", "Statistics", "Research Methods"]
    return ["Scientific Method", "Logic", "Mathematics", "Python"]

# --- SIDEBAR SETTINGS ---
st.sidebar.header("🔍 Global Search Settings")
topics = ["Artificial Intelligence", "Astrophysics", "Bio-medical Engineering", "Black Holes", 
          "Climate Change", "Deep Learning", "Data Science", "Genetics", "Robotics", "--- TYPE MY OWN ---"]

selection = st.sidebar.selectbox("Choose a Field:", options=topics, index=None, placeholder="Search...")
user_query = st.sidebar.text_input("Custom Topic:") if selection == "--- TYPE MY OWN ---" else (selection if selection else "Research")

blob = TextBlob(user_query)
query = str(blob.correct())
if query.lower() != user_query.lower():
    st.sidebar.warning(f"Searching for: **{query}**")

num_papers = st.sidebar.slider("Papers to analyze", 5, 20, 10)
simplify = st.sidebar.toggle("Enable ELI15 Mode (Simple Language)")

# --- MAIN TABS ---
tab1, tab2 = st.tabs(["🌐 Global ArXiv Search", "📂 Personal Document Analysis"])

with tab1:
    if st.sidebar.button("Launch Global Analysis 🚀"):
        with st.spinner("Fetching latest data..."):
            search = arxiv.Search(query=query, max_results=num_papers, sort_by=arxiv.SortCriterion.Relevance)
            results = [{"Title": r.title, "Summary": r.summary, "URL": r.pdf_url, "Date": r.published.date()} for r in search.results()]
            
            if results:
                df = pd.DataFrame(results)
                c1, c2 = st.columns([2, 1])
                with c1:
                    st.subheader("🔥 Trends Cloud")
                    wc = WordCloud(width=800, height=400, background_color="black", colormap="magma").generate(" ".join(df['Summary']))
                    fig, ax = plt.subplots(); ax.imshow(wc); ax.axis("off"); st.pyplot(fig)
                with c2:
                    st.subheader("🎓 Career Path")
                    if GRAPHVIZ_AVAILABLE:
                        dot = graphviz.Digraph(); dot.attr(rankdir='LR')
                        sk = get_career_advice(query)
                        dot.node('A', 'School'); dot.node('B', sk[0]); dot.node('C', sk[1]); dot.node('D', 'Expert')
                        dot.edges(['AB', 'BC', 'CD']); st.graphviz_chart(dot)
                
                st.divider()
                for _, row in df.iterrows():
                    with st.expander(f"📌 {row['Title']}"):
                        if simplify: st.info(f"**Simplified:** {generate_ai_summary(row['Summary'], 1)}")
                        else: st.write(row['Summary'])
                        st.link_button("PDF Link", row['URL'])
            else: st.error("No papers found.")

with tab2:
    st.subheader("Analyze Your PDF, DOCX, or PPTX")
    up_file = st.file_uploader("Upload file", type=["pdf", "docx", "pptx"])
    if up_file:
        raw_text = extract_text_from_file(up_file)
        if raw_text:
            st.success(f"File '{up_file.name}' loaded!")
            sc1, sc2 = st.columns([2, 1])
            with sc1:
                st.subheader("📊 Document WordCloud")
                wc2 = WordCloud(width=800, height=400, background_color="white").generate(raw_text)
                fig2, ax2 = plt.subplots(); ax2.imshow(wc2); ax2.axis("off"); st.pyplot(fig2)
